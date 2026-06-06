from __future__ import annotations

from pathlib import Path
from textwrap import dedent

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


OUT = Path(__file__).resolve().parent
BASE = "agentic_pm_substrate_thesis_2026-06-06"
TITLE = "Design the Interactions First"
SUBTITLE = "pm-substrate and the Project-Management Layer for Agentic Workspaces"
AUTHOR = "Emmanuel Akinwale, JOAT Labs"
DATE = "June 6, 2026"


SECTIONS = [
    {
        "heading": "Abstract",
        "body": [
            "Modern organizations do not suffer from a shortage of software. They suffer from a shortage of coherent operational state. Businesses already use capable tools for sales, finance, planning, communication, documents, analytics, and delivery. The breakdown happens between those tools: state is copied, context is lost, ownership is unclear, and each system maintains a partial local model of work.",
            "The arrival of AI agents makes this problem more urgent. Large language models can generate plans and call tools, but a model call is not a system of record, a prompt is not a workspace, and retrieved memory is not authority. Agents are bounded perception-action systems. They act from partial observations, compressed representations, uncertain beliefs, and delayed feedback. Without a governed operational substrate, agents inherit the same fragmentation that already makes modern SaaS work brittle.",
            "This thesis argues that pm-substrate is the project-management layer for agentic workspaces. It formalizes the operational-state function of project management: maintaining shared goals, valid next actions, owners, blockers, dependencies, evidence, authority, and continuity across humans, tools, workflows, and AI agents. Its central design law is simple: no agent's local state is authoritative unless it is reconciled with substrate state.",
        ],
    },
    {
        "heading": "1. The Strategic Problem: The System Is the Strategy",
        "body": [
            "The modern workspace has been optimized around components. Every function has a specialized application, and many of those applications are excellent inside their own boundaries. Yet the operating experience is still fragmented because the value of work does not live inside any one component. It lives in the interactions between components.",
            "This is the systems insight behind JOAT Labs. A business is not a CRM, a project tracker, a document store, an inbox, a calendar, and a finance system placed beside one another. A business is the coordinated behavior that emerges when those capabilities exchange meaning at the right time, under the right authority, with the right permissions, and in the right workflow position.",
            "The dominant software response has been to add more interfaces, more integrations, and now more AI features to existing silos. That response treats the workspace as a shelf of tools. pm-substrate treats the workspace as a system of bounded actors whose local models must be reconciled into shared operational state.",
        ],
    },
    {
        "heading": "2. Project Management as Institutional State Governance",
        "body": [
            "A project manager is often described as the person who schedules meetings, tracks tasks, and follows up with owners. That is too small. At a deeper level, project management is the organizational discipline for maintaining shared state under partial observation and change.",
            "A project manager keeps the current objective visible, translates between specialized teams, records decisions, watches dependencies, detects blockers, manages risks, clarifies ownership, and determines which next actions are valid. These are all state functions. They answer questions such as: What is true now? Who owns it? Which source is binding? What changed? Which step is next? What is blocked? What evidence supports the decision?",
            "This is why the project-management layer is the correct metaphor for the modern workspace. The tools a business uses are like departments in a cross-functional team. Each department can be excellent and still fail as a whole if their state does not synchronize. pm-substrate is the missing project manager for tools, workflows, people, and agents.",
        ],
    },
    {
        "heading": "3. What an Agent Is, Structurally",
        "body": [
            "An agent is not first a chatbot or an automation. At the lowest useful level, an agent is a bounded system coupled to an environment through perception, internal representation, action, and feedback.",
            "The environment has some state. The agent receives partial observations from that environment. The agent maintains an internal model or belief. It updates that model as observations arrive. It selects an action according to a goal, policy, or objective. The action changes the environment or the agent's future observations. Feedback then arrives, sometimes immediately and sometimes late.",
            "This structure exists before artificial intelligence. People, departments, software services, robots, and LLM-based agents all fit the same pattern. The agent-state problem is therefore not a narrow LLM implementation issue. It is the old coordination problem of bounded actors acting from incomplete local models.",
        ],
    },
    {
        "heading": "4. From Numbers to Agentic State",
        "body": [
            "A large language model begins as numbers: token IDs, embeddings, tensors, matrix operations, and learned parameters. Statistical learning turns those numbers into a function that predicts outputs from inputs. A transformer language model turns prior tokens into a probability distribution over next tokens.",
            "That mechanism is powerful, but it is not operational state. Model weights are parametric state: a compressed record of historical training and fine-tuning data. Prompt context is inference state: a temporary working window for one call or run. Retrieval and memory are retrieval state: external records brought into context, often without full lifecycle governance. None of these is automatically current, authoritative, complete, or safe to mutate against.",
            "The agent-state failure begins when a statistical predictor is promoted into an actor without giving it a governed current-state layer between prediction and action. The model can reason well and still act badly if its context is stale, missing, conflicting, or non-authoritative.",
        ],
    },
    {
        "heading": "5. State Taxonomy",
        "body": [
            "The word state is overloaded, so the thesis must separate its meanings. Reality state is what is actually true in the environment. Observed state is what an actor can see. Represented state is how observations are encoded. Belief state is what the actor thinks is true. Operational state is the current actionable configuration of work. Authoritative state is the state the institution accepts as binding. Historical state is the event sequence that explains how current state came to be. Projected state is a derived view built from authoritative records.",
            "pm-substrate does not claim to govern all of reality. It governs operational, authoritative, historical, and projected workspace state. That distinction matters. The substrate is not the agent's entire world model. It is the governed operating reality against which work actions are validated.",
        ],
    },
    {
        "heading": "6. The Core Thesis",
        "body": [
            "Modern workspace failure is a multi-agent state coherence problem. Each human, tool, workflow, and AI agent acts from a local model. Failures occur when those local models diverge and no shared substrate determines which facts are current, which sources are authoritative, and which actions are valid.",
            "Project management is the human discipline that has historically handled this problem. It maintains institutional state across specialists. pm-substrate formalizes that project-management state function in software for agentic workspaces.",
            "The shortest version is: pm-substrate is not a bigger memory for agents. It is the missing state-estimation, authority, and project-management layer between statistical prediction and operational action.",
        ],
    },
    {
        "heading": "7. What pm-substrate Is",
        "body": [
            "pm-substrate is a tenant-scoped operational substrate for coordinating work across humans, tools, workflows, and AI agents. It is not a point-to-point integration script and not a single vertical application. It is the shared state layer that allows specialized systems to coordinate through declared contracts instead of private assumptions.",
            "Its minimal primitive set includes actors, entities, relationships, events, capabilities, policies, workflow state, evidence, projections, commitments, and continuity records. Together, these primitives turn scattered observations into governed operational state.",
            "The design law is: no agent's local state is authoritative unless it is reconciled with substrate state. Agent memory can propose. Agent belief can explain. Agent plans can recommend. But authoritative project state changes only through substrate transitions.",
        ],
    },
    {
        "heading": "8. Authority Promotion: From Observation to Operational Reality",
        "body": [
            "A substrate needs a promotion ladder. Otherwise, shared state becomes a slogan rather than a control boundary.",
            "The ladder is: observation, evidence, proposed fact, validated fact, authoritative transition, projection update, and continuity update. A human, tool, or agent observes something. The observation is recorded with source, time, actor, and provenance. The system interprets the evidence as a candidate statement about an entity or workflow. Schema, permission, freshness, and conflict checks pass. An event is appended to the log. Derived views and workflow states are recomputed. Unresolved work, decisions, and evidence are preserved for future agents.",
            "This is the route from local perception to institutional state. It is also the route that prevents a model's confident output from becoming an unvalidated mutation.",
        ],
    },
    {
        "heading": "9. Agent Action Lifecycle",
        "body": [
            "In pm-substrate, the LLM is a proposal engine. It can interpret, summarize, plan, and recommend. But a proposed action is not the same as an authorized mutation.",
            "The action lifecycle is: build a current_state_view, ask the model for a proposal, record the proposal's read set and observation contract, validate the proposal against current substrate state, warn or block based on enforcement mode, execute only through a capability boundary, and append evidence-backed events for what happened.",
            "This lifecycle turns agent behavior into project-management behavior. The agent must show what it relied on, whether those reads are still fresh, whether the workflow position still permits the action, whether the authority rule matches, and whether the action subject matches the current state view.",
        ],
    },
    {
        "heading": "10. Evaluation Strategy",
        "body": [
            "The thesis must be falsifiable. If pm-substrate does real work, agents should complete changing cross-tool workflows with lower stale-action rate, lower state-disagreement rate, lower rework, and higher evidence coverage than agents using chat history plus raw tool access.",
            "The current evaluation families are stale observation, source authority conflict, workflow invalidation, representation loss, memory drift, capability contract violation, and continuity break. The implementation deliberately separates scaffolded scenarios from stronger evidence stages: detected warning, blocked mutation, and paired behavioral improvement.",
            "That evidence maturity ladder matters. A scenario that merely says the substrate should pass is not proof. A system that emits a deterministic warning is stronger. A mutation gate that blocks a bad action is stronger still. A paired behavioral improvement in an executable workflow is the strongest version for the current project.",
        ],
    },
    {
        "heading": "11. Current Implementation Evidence",
        "body": [
            "The repository already implements the first concrete spine of the thesis. It includes typed events, graph state, workflow state, projection state, continuity concepts, capability contracts, eval metrics, and an agent-state package.",
            "The agent-state implementation defines CurrentStateView, StateRef, ReadSetEntry, ProposedAction, ObservationContract, and ActionProposalReview. It can warn when read-set refs are stale, required sources are missing, authority differs, projection versions changed, workflow positions conflict, current views contain conflicts, or a proposed action targets a subject different from the current-state view.",
            "The ArrowHedge adapter gives the thesis a high-consequence tool surface. Source records are parsed, validated against profile/entity mapping, emitted as typed events, folded into a Common Operating Picture, transformed into a current-state view, and reviewed before action. That path tests the actual substrate strategy: source records -> semantic mapping -> typed events -> graph/projection -> current-state view -> proposal review -> eval metrics.",
        ],
    },
    {
        "heading": "12. Why This Is Not iPaaS, Workflow Automation, or RAG",
        "body": [
            "iPaaS and workflow automation usually begin from connections: when this happens in one app, do that in another. pm-substrate begins from shared operational state. It asks which entity this is about, which source owns the fact, whether the downstream action is still valid, what changed since the trigger fired, and what history can be replayed.",
            "RAG is also insufficient by itself. Retrieval can bring information into context, but retrieval is not governance. It does not automatically provide authority, freshness, workflow validity, permissions, contradiction handling, or mutation safety.",
            "The same is true of larger context windows. More context increases capacity, not authority. A bigger prompt can still contain stale, conflicting, or non-binding information. pm-substrate exists because the problem is not just memory. The problem is governed operational state.",
        ],
    },
    {
        "heading": "13. Strategic Implications for JOAT Labs",
        "body": [
            "The JOAT Labs strategy is not to compete feature-by-feature with existing tools. The strategy is to own the interaction layer that makes tools and agents behave as one coherent workspace. The system is the strategy.",
            "The substrate becomes the durable center of the product: not because every user wants to see a database or event log, but because every reliable agentic workflow depends on the substrate's guarantees. Users experience it as fewer repeated explanations, fewer stale handoffs, clearer owners, stronger continuity, and safer automation.",
            "The product wedge should therefore be a visible Common Operating Picture plus validation gates for real workflows. The deep architecture can be sophisticated, but the user-facing promise is simple: many bounded actors can work from one coherent operating reality.",
        ],
    },
    {
        "heading": "14. Roadmap",
        "body": [
            "The next phase is to move from warning artifacts into stronger execution evidence. First, expand the Common Operating Picture so it becomes the primary state surface for humans and agents. Second, wire advisory proposal reviews into a real workflow or capability mutation path. Third, add blocking-mode gates for selected high-consequence transitions. Fourth, run paired evals that measure behavior, not just detection. Fifth, connect continuity records so agent resumes are evidence-linked rather than transcript-dependent.",
            "The long-term research program should continue to combine project-management research, multi-agent systems, state estimation, distributed systems, workflow theory, and LLM agent evaluation. The substrate thesis is strongest when it refuses to live in only one discipline.",
        ],
    },
    {
        "heading": "Conclusion",
        "body": [
            "The next era of workspace software will not be won by adding isolated AI features to isolated tools. That path multiplies fragmented state. The next era will be won by systems that make humans, tools, workflows, and agents interoperable at the level where work actually happens: entities, events, authority, permissions, workflows, evidence, commitments, and continuity.",
            "An agent should not be treated as a magic worker with memory. It is a bounded perception-action system operating under partial observation. It needs a governed operational substrate to act safely over time.",
            "That is the role of pm-substrate. It is the project-management layer for agentic workspaces: the system that keeps bounded actors aligned around shared state while preserving authority, causality, permission, and accountability. The principle is simple: design the interactions first.",
        ],
    },
]


TABLES = [
    (
        "State Layers in Agentic Work",
        ["State type", "Where it lives", "Why it is insufficient"],
        [
            ["Parametric state", "Model weights", "Historical statistical compression, not current authority."],
            ["Inference state", "Prompt, activations, context window", "Temporary, finite, order-sensitive, and not durable."],
            ["Retrieval state", "Documents, vector memory, summaries", "Can be stale, conflicting, or missing source authority."],
            ["Operational state", "Events, graph, workflows, projections", "Only useful when governed, validated, and current."],
        ],
    ),
    (
        "pm-substrate Primitive Map",
        ["Agent-state need", "pm-substrate primitive"],
        [
            ["Environment state", "Entity graph"],
            ["State changes", "Append-only event log"],
            ["Possible actions", "Capability registry"],
            ["Valid process", "Workflow runtime"],
            ["Allowed actions", "Policies and validation rules"],
            ["Observed facts", "Evidence-linked projections"],
            ["Agent continuity", "Continuity ledger"],
            ["Feedback", "Workflow outcomes and emitted events"],
        ],
    ),
    (
        "Evidence Maturity Ladder",
        ["Stage", "Meaning"],
        [
            ["Scaffolded scenario", "A designed case that explains the failure class but does not prove behavior."],
            ["Detected warning", "The substrate deterministically surfaces stale, missing, conflicted, or non-authoritative state."],
            ["Blocked mutation", "A gate prevents a bad transition from mutating project state."],
            ["Paired behavioral improvement", "The substrate arm outperforms a baseline in an executable paired workflow."],
        ],
    ),
]


REFERENCES = [
    "Ackoff, R. L. On Systems Thinking. W. Edwards Deming Institute summary.",
    "Bertalanffy, L. von. General System Theory: Foundations, Development, Applications. 1968.",
    "Brown, T. et al. Language Models are Few-Shot Learners. 2020.",
    "Fisher, D. A. An Emergent Perspective on Interoperation in Systems of Systems. SEI/CMU, 2006.",
    "Kaelbling, L. P., Littman, M. L., and Cassandra, A. R. Planning and acting in partially observable stochastic domains. 1998.",
    "Kalman, R. E. A New Approach to Linear Filtering and Prediction Problems. 1960.",
    "Lewis, P. et al. Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks. 2020.",
    "Maier, M. W. Architecting Principles for Systems-of-Systems. 1998.",
    "Meadows, D. H. Thinking in Systems: A Primer. 2008.",
    "MuleSoft / Salesforce. Connectivity Benchmark Reports, 2025-2026.",
    "Park, J. S. et al. Generative Agents: Interactive Simulacra of Human Behavior. 2023.",
    "Senge, P. M. The Fifth Discipline. 1990 / 2006.",
    "Sutton, R. S. and Barto, A. G. Reinforcement Learning: An Introduction. 2020.",
    "Trist, E. L. and Bamforth, K. W. Some Social and Psychological Consequences of the Longwall Method of Coal-Getting. 1951.",
    "Vaswani, A. et al. Attention Is All You Need. 2017.",
    "Yao, S. et al. ReAct: Synergizing Reasoning and Acting in Language Models. 2023.",
    "JOAT Labs source documents: The System Is the Strategy; Design the Interactions First; Agent From Numbers To State.",
]


SLIDES = [
    ("Design the Interactions First", ["pm-substrate and the project-management layer for agentic workspaces", "Repo thesis deck | June 2026"]),
    ("One-Sentence Thesis", ["Modern workspace failure is a shared-state problem.", "AI agents make it urgent because they can act without knowing what is current, authoritative, or valid.", "pm-substrate supplies the project-management layer between prediction and action."]),
    ("The Workspace Problem", ["Components are mature; interactions are brittle.", "Each tool maintains a local model of work.", "Humans become the manual integration layer.", "AI bolt-ons multiply the number of local models unless state is governed."]),
    ("Project Management Is State Governance", ["PM maintains the operating reality of a project: goals, owners, blockers, risks, decisions, dependencies, and next valid actions.", "That function is not just administrative. It is institutional state management."]),
    ("What an Agent Is", ["Agent = bounded perception-action loop.", "It observes partially, maintains a local model, proposes actions, acts, and receives feedback.", "The agent-state problem is older than LLMs; LLMs expose it at scale."]),
    ("The LLM State Gap", ["Weights are parametric state.", "Prompts are inference state.", "Memory/RAG is retrieval state.", "None of these is automatically operational authority."]),
    ("pm-substrate Design Law", ["No agent's local state is authoritative unless it is reconciled with substrate state.", "Agents can propose. The substrate validates. Mutations happen only through governed transitions."]),
    ("Core Primitives", ["Actors", "Entities and relationships", "Append-only events", "Capabilities and policies", "Workflow state", "Evidence and projections", "Commitments and continuity records"]),
    ("Authority Promotion Ladder", ["Observation", "Evidence", "Proposed fact", "Validated fact", "Authoritative transition", "Projection update", "Continuity update"]),
    ("Proposal Review Lifecycle", ["Build current_state_view.", "Capture read set and observation contract.", "Validate freshness, authority, workflow position, projection version, conflicts, and subject.", "Warn-first today; blocking gates next."]),
    ("Evaluation Strategy", ["Failure classes: stale observation, source authority conflict, workflow invalidation, representation loss, memory drift, capability contract violation.", "Evidence stages: scaffolded scenario, detected warning, blocked mutation, paired behavioral improvement."]),
    ("Current Repo Evidence", ["@pm/agent-state contracts", "ArrowHedge source records -> typed events -> COP -> current-state view", "Read-set and proposal review warnings", "Eval maturity labels to avoid overclaiming proof"]),
    ("Strategic Wedge", ["A visible Common Operating Picture for humans and agents.", "Validation gates for high-consequence workflows.", "Continuity records that let agents resume from evidence rather than transcript memory."]),
    ("Closing Claim", ["pm-substrate is not a bigger memory for agents.", "It is the missing state-estimation, authority, and project-management layer between statistical prediction and operational action."]),
]


def markdown_text() -> str:
    parts = [
        f"# {TITLE}: {SUBTITLE}",
        "",
        f"**Author:** {AUTHOR}  ",
        f"**Date:** {DATE}  ",
        "**Status:** Unified repo thesis paper",
        "",
    ]
    for section in SECTIONS:
        parts.append(f"## {section['heading']}")
        parts.append("")
        for paragraph in section["body"]:
            parts.append(paragraph)
            parts.append("")
        if section["heading"] == "5. State Taxonomy":
            parts.extend(markdown_table(TABLES[0]))
        if section["heading"] == "7. What pm-substrate Is":
            parts.extend(markdown_table(TABLES[1]))
        if section["heading"] == "10. Evaluation Strategy":
            parts.extend(markdown_table(TABLES[2]))
    parts.append("## References")
    parts.append("")
    for ref in REFERENCES:
        parts.append(f"- {ref}")
    parts.append("")
    return "\n".join(parts)


def markdown_table(table_spec):
    title, headers, rows = table_spec
    lines = [f"**{title}**", ""]
    lines.append("| " + " | ".join(headers) + " |")
    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
    for row in rows:
        lines.append("| " + " | ".join(row) + " |")
    lines.append("")
    return lines


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def add_docx_table(doc: Document, table_spec) -> None:
    title, headers, rows = table_spec
    p = doc.add_paragraph()
    run = p.add_run(title)
    run.bold = True
    run.font.color.rgb = RGBColor(31, 78, 121)
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        set_cell_shading(cell, "D9EAF7")
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = value
            cells[i].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    doc.add_paragraph()


def build_docx(path: Path) -> None:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.8)
    styles["Normal"].paragraph_format.space_after = Pt(6)
    for style_name, size, color in [
        ("Heading 1", 18, RGBColor(31, 78, 121)),
        ("Heading 2", 14, RGBColor(31, 78, 121)),
        ("Heading 3", 12, RGBColor(54, 95, 145)),
    ]:
        style = styles[style_name]
        style.font.name = "Aptos Display"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = color

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title.add_run(TITLE)
    title_run.bold = True
    title_run.font.size = Pt(26)
    title_run.font.color.rgb = RGBColor(31, 78, 121)
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = subtitle.add_run(SUBTITLE)
    sub_run.font.size = Pt(14)
    sub_run.italic = True
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"{AUTHOR} | {DATE}").font.size = Pt(10)
    doc.add_paragraph()

    callout = doc.add_paragraph()
    callout.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = callout.add_run("Thesis: pm-substrate is the missing project-management layer between statistical prediction and operational action.")
    r.bold = True
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(77, 77, 77)

    doc.add_page_break()
    for section_data in SECTIONS:
        doc.add_heading(section_data["heading"], level=1 if section_data["heading"] in ["Abstract", "Conclusion"] else 2)
        for paragraph in section_data["body"]:
            doc.add_paragraph(paragraph)
        if section_data["heading"] == "5. State Taxonomy":
            add_docx_table(doc, TABLES[0])
        if section_data["heading"] == "7. What pm-substrate Is":
            add_docx_table(doc, TABLES[1])
        if section_data["heading"] == "10. Evaluation Strategy":
            add_docx_table(doc, TABLES[2])

    doc.add_heading("References", level=1)
    for ref in REFERENCES:
        p = doc.add_paragraph(style=None)
        p.style = doc.styles["Normal"]
        p.paragraph_format.left_indent = Inches(0.18)
        p.add_run(ref)

    footer = doc.sections[0].footer.paragraphs[0]
    footer.text = "JOAT Labs | pm-substrate unified thesis"
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.save(path)


def pdf_styles():
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name="TitleMain",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=25,
        leading=30,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#1F4E79"),
        spaceAfter=8,
    ))
    styles.add(ParagraphStyle(
        name="Subtitle",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=13,
        leading=17,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#444444"),
        spaceAfter=18,
    ))
    styles.add(ParagraphStyle(
        name="H1Custom",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=15,
        leading=18,
        textColor=colors.HexColor("#1F4E79"),
        spaceBefore=12,
        spaceAfter=6,
    ))
    styles.add(ParagraphStyle(
        name="BodyCustom",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=9.8,
        leading=14,
        spaceAfter=7,
    ))
    styles.add(ParagraphStyle(
        name="Ref",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=8.5,
        leading=11,
        leftIndent=8,
        spaceAfter=4,
    ))
    return styles


def add_pdf_table(story, table_spec, styles) -> None:
    title, headers, rows = table_spec
    story.append(Paragraph(f"<b>{title}</b>", styles["BodyCustom"]))
    data = [headers] + rows
    col_count = len(headers)
    col_widths = [6.9 * inch / col_count] * col_count
    table = Table(data, colWidths=col_widths, hAlign="LEFT", repeatRows=1)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAF7")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#17365D")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 8.2),
        ("LEADING", (0, 0), (-1, -1), 10),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#9BBAD4")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(table)
    story.append(Spacer(1, 0.14 * inch))


def build_pdf(path: Path) -> None:
    styles = pdf_styles()
    doc = SimpleDocTemplate(
        str(path),
        pagesize=letter,
        rightMargin=0.72 * inch,
        leftMargin=0.72 * inch,
        topMargin=0.68 * inch,
        bottomMargin=0.62 * inch,
    )
    story = [
        Paragraph(TITLE, styles["TitleMain"]),
        Paragraph(SUBTITLE, styles["Subtitle"]),
        Paragraph(f"{AUTHOR}<br/>{DATE}<br/><br/><b>Unified repo thesis paper</b>", styles["Subtitle"]),
        Spacer(1, 0.3 * inch),
        Paragraph("<b>Thesis:</b> pm-substrate is the missing project-management layer between statistical prediction and operational action.", styles["BodyCustom"]),
        PageBreak(),
    ]

    for section_data in SECTIONS:
        story.append(Paragraph(section_data["heading"], styles["H1Custom"]))
        for paragraph in section_data["body"]:
            story.append(Paragraph(paragraph, styles["BodyCustom"]))
        if section_data["heading"] == "5. State Taxonomy":
            add_pdf_table(story, TABLES[0], styles)
        if section_data["heading"] == "7. What pm-substrate Is":
            add_pdf_table(story, TABLES[1], styles)
        if section_data["heading"] == "10. Evaluation Strategy":
            add_pdf_table(story, TABLES[2], styles)

    story.append(Paragraph("References", styles["H1Custom"]))
    for ref in REFERENCES:
        story.append(Paragraph(ref, styles["Ref"]))

    def footer(canvas, doc_obj):
        canvas.saveState()
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(colors.HexColor("#777777"))
        canvas.drawString(0.72 * inch, 0.38 * inch, "JOAT Labs | pm-substrate unified thesis")
        canvas.drawRightString(7.78 * inch, 0.38 * inch, f"Page {doc_obj.page}")
        canvas.restoreState()

    doc.build(story, onFirstPage=footer, onLaterPages=footer)


def add_textbox(slide, x, y, w, h, text, size=24, bold=False, color="FFFFFF", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color)
    return box


def add_bullets(slide, x, y, w, h, bullets, size=18):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PptPt(size)
        p.font.color.rgb = PptRGBColor(35, 45, 55)
        p.space_after = PptPt(8)
    return box


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    for idx, (title, bullets) in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PptRGBColor(246, 248, 250)

        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PptInches(13.333), PptInches(0.22))
        band.fill.solid()
        band.fill.fore_color.rgb = PptRGBColor(31, 78, 121)
        band.line.fill.background()

        add_textbox(slide, 0.62, 0.55, 11.7, 0.68, title, size=27 if idx else 34, bold=True, color="1F4E79")
        if idx == 0:
            add_textbox(slide, 0.68, 1.45, 11.8, 0.7, bullets[0], size=22, color="334155")
            add_textbox(slide, 0.68, 2.3, 6.2, 0.4, bullets[1], size=15, color="64748B")
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(0.7), PptInches(4.75), PptInches(11.9), PptInches(1.1))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptRGBColor(31, 78, 121)
            shape.line.fill.background()
            add_textbox(slide, 1.05, 5.03, 11.2, 0.5, "The system is the strategy. The substrate is the PM layer.", size=22, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
        else:
            add_bullets(slide, 0.85, 1.55, 8.3, 4.9, bullets, size=19 if len(bullets) <= 3 else 17)
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(9.55), PptInches(1.58), PptInches(2.95), PptInches(4.85))
            panel.fill.solid()
            panel.fill.fore_color.rgb = PptRGBColor(232, 242, 249)
            panel.line.color.rgb = PptRGBColor(155, 186, 212)
            if idx in [4, 5, 6, 8, 9]:
                label = ["Observe", "Represent", "Validate", "Act"][idx % 4]
            elif idx in [10, 11]:
                label = "Measure"
            else:
                label = "Coordinate"
            add_textbox(slide, 9.85, 3.34, 2.35, 0.55, label, size=24, bold=True, color="1F4E79", align=PP_ALIGN.CENTER)
        add_textbox(slide, 0.65, 7.05, 5.6, 0.24, "JOAT Labs | pm-substrate thesis", size=8.5, color="64748B")
        add_textbox(slide, 11.65, 7.05, 1.0, 0.24, str(idx + 1), size=8.5, color="64748B", align=PP_ALIGN.RIGHT)

    prs.save(path)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    md_path = OUT / f"{BASE}.md"
    docx_path = OUT / f"{BASE}.docx"
    pdf_path = OUT / f"{BASE}.pdf"
    pptx_path = OUT / f"{BASE}_deck.pptx"

    md_path.write_text(markdown_text(), encoding="utf-8")
    build_docx(docx_path)
    build_pdf(pdf_path)
    build_pptx(pptx_path)
    print(md_path)
    print(docx_path)
    print(pdf_path)
    print(pptx_path)


if __name__ == "__main__":
    main()
